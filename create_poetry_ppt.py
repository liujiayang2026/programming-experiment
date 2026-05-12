"""Generate poetry recitation background PPT using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import random, math

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
TOTAL = 7


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, r, g, b):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(r, g, b)


def tb(slide, l, t, w, h, txt, sz=18, bold=False, color=RGBColor(0xFF,0xFF,0xFF),
       align=PP_ALIGN.CENTER, font='Arial'):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.font.name = font
    return tf


def multi_tb(slide, l, t, w, h, lines, sz=18, color=RGBColor(0xFF,0xFF,0xFF),
             align=PP_ALIGN.CENTER, font='Arial', spacing=1.4):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.alignment = align
        p.font.name = font
        p.space_after = Pt(int(sz * (spacing - 1)))
    return tf


def accent(slide, l, t, w, r, g, b):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(r, g, b)
    line.line.fill.background()


def page_no(slide, n):
    tb(slide, 11.8, 7.0, 1.3, 0.35, f"{n}/{TOTAL}",
       sz=9, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.RIGHT)


def circle(slide, cx, cy, r, rc, gc, bc):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(cx - r), Inches(cy - r), Inches(r * 2), Inches(r * 2))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(rc, gc, bc)
    s.line.fill.background()


def blend(c1, c2, ratio):
    """Blend two colors. ratio=0 means pure c1, ratio=1 means pure c2."""
    return (
        int(c1[0] + (c2[0] - c1[0]) * ratio),
        int(c1[1] + (c2[1] - c1[1]) * ratio),
        int(c1[2] + (c2[2] - c1[2]) * ratio),
    )


# ================================================================
#  SLIDE 1 - Opening: Lonely Cloud
# ================================================================
s1 = blank_slide()
bg1 = (0x5B, 0x8C, 0xB8)
set_bg(s1, *bg1)

circle(s1, 3.5, 2.8, 1.8, *blend(bg1, (0xFF,0xFF,0xFF), 0.25))
circle(s1, 5.2, 2.2, 1.3, *blend(bg1, (0xFF,0xFF,0xFF), 0.22))
circle(s1, 4.3, 3.5, 1.0, *blend(bg1, (0xFF,0xFF,0xFF), 0.18))

multi_tb(s1, 1.5, 1.5, 10.3, 1.8,
    ["I wandered lonely as a cloud"],
    sz=42, color=RGBColor(0xFF,0xFF,0xFF), font='Arial')

accent(s1, 5.0, 3.6, 3.3, 0xFF, 0xD7, 0x00)

multi_tb(s1, 1.5, 3.9, 10.3, 1.2,
    ["\u4E91  \u00B7  \u5B64  \u00B7  \u98D8"],
    sz=36, color=RGBColor(0xE8,0xE8,0xE8), font='Microsoft YaHei')

multi_tb(s1, 2.0, 5.3, 9.3, 0.8,
    ["\u4E00\u6735\u5B64\u4E91\u98D8\u8FC7\u9AD8\u9AD8\u7684\u5C71\u8C37"],
    sz=16, color=RGBColor(0xCC,0xDD,0xEE), font='Microsoft YaHei')

tb(s1, 1.0, 6.4, 11.3, 0.5, "\u2014 William Wordsworth",
   sz=14, color=RGBColor(0xAA,0xBB,0xCC), font='Arial')

page_no(s1, 1)

# ================================================================
#  SLIDE 2 - Encounter: Golden Daffodils
# ================================================================
s2 = blank_slide()
bg2 = (0x2E, 0x7D, 0x32)
set_bg(s2, *bg2)

for angle in range(0, 360, 60):
    rad = angle * math.pi / 180
    circle(s2, 9.3 + 1.5 * math.cos(rad), 3.5 + 1.5 * math.sin(rad),
           0.6, *blend(bg2, (0xFF,0xD7,0x00), 0.35))

multi_tb(s2, 0.8, 1.3, 8.0, 2.6,
    ["a crowd,", "a host of golden daffodils"],
    sz=38, color=RGBColor(0xFF,0xFF,0xFF), font='Arial')

accent(s2, 2.8, 4.0, 3.5, 0xFF, 0xD7, 0x00)

multi_tb(s2, 0.8, 4.3, 8.0, 1.4,
    ["\u91D1\u8272  \u00B7  \u6210\u7FA4  \u00B7  \u821E\u52A8"],
    sz=34, color=RGBColor(0xFF,0xEE,0x88), font='Microsoft YaHei')

multi_tb(s2, 1.2, 5.5, 7.5, 0.8,
    ["\u6E56\u8FB9\u6811\u4E0B\uFF0C\u4E00\u5927\u7247\u91D1\u8272\u6C34\u4ED9\u5728\u98CE\u4E2D\u821E\u52A8"],
    sz=16, color=RGBColor(0xCC,0xFF,0xCC), font='Microsoft YaHei')

page_no(s2, 2)

# ================================================================
#  SLIDE 3 - Stretching: Like Milky Way Stars
# ================================================================
s3 = blank_slide()
bg3 = (0x0D, 0x1B, 0x2A)
set_bg(s3, *bg3)

for _ in range(40):
    sx = random.uniform(0.3, 13.0)
    sy = random.uniform(0.2, 6.8)
    sr = random.uniform(0.02, 0.06)
    br = random.randint(160, 255)
    circle(s3, sx, sy, sr, br, br, br)

multi_tb(s3, 1.0, 1.2, 11.3, 2.8,
    ["continuous as the stars", "never-ending line"],
    sz=38, color=RGBColor(0xFF,0xFF,0xFF), font='Arial')

accent(s3, 5.0, 4.3, 3.3, 0xFF, 0xD7, 0x00)

multi_tb(s3, 1.0, 4.6, 11.3, 1.3,
    ["\u7E41\u661F  \u00B7  \u7EF5\u5EF6  \u00B7  \u6D77\u6E7E"],
    sz=34, color=RGBColor(0xFF,0xEE,0xAA), font='Microsoft YaHei')

multi_tb(s3, 1.5, 5.5, 10.3, 0.8,
    ["\u6C34\u4ED9\u6CBF\u7740\u6D77\u6E7E\u65E0\u5C3D\u7EF5\u5EF6\uFF0C\u50CF\u94F6\u6CB3\u4E2D\u7684\u7E41\u661F"],
    sz=16, color=RGBColor(0xAA,0xCC,0xEE), font='Microsoft YaHei')

page_no(s3, 3)

# ================================================================
#  SLIDE 4 - Joyful: Outdoing the Waves
# ================================================================
s4 = blank_slide()
bg4 = (0x1A, 0x5C, 0x8A)
set_bg(s4, *bg4)

for i in range(0, 15):
    circle(s4, i * 0.95 - 0.1, 5.5 + 0.4 * math.sin(i * 0.75),
           0.45, *blend(bg4, (0x34,0x98,0xDB), 0.55))

multi_tb(s4, 0.8, 1.1, 11.5, 3.0,
    ["tossing their heads in sprightly dance",
     "out-did the sparkling waves"],
    sz=34, color=RGBColor(0xFF,0xFF,0xFF), font='Arial')

accent(s4, 5.0, 4.2, 3.3, 0xFF, 0xD7, 0x00)

multi_tb(s4, 0.8, 4.5, 11.5, 1.2,
    ["\u96C0\u8DC3  \u00B7  \u80DC\u8FC7\u6D6A\u82B1"],
    sz=34, color=RGBColor(0xFF,0xEE,0xBB), font='Microsoft YaHei')

multi_tb(s4, 2.0, 5.2, 9.3, 0.8,
    ["\u6C34\u4ED9\u6602\u5934\u96C0\u8DC3\u8D77\u821E\uFF0C\u8FDE\u95EA\u4EAE\u7684\u6CE2\u6D6A\u90FD\u6BD4\u4E0D\u8FC7"],
    sz=16, color=RGBColor(0xCC,0xEE,0xFF), font='Microsoft YaHei')

page_no(s4, 4)

# ================================================================
#  SLIDE 5 - Gazing: Invisible Wealth
# ================================================================
s5 = blank_slide()
bg5 = (0x5D, 0x40, 0x37)
set_bg(s5, *bg5)

circle(s5, 9.0, 2.5, 2.8, *blend(bg5, (0xFF,0xD7,0x00), 0.18))
circle(s5, 9.0, 2.5, 2.0, *blend(bg5, (0xFF,0xC1,0x07), 0.25))

multi_tb(s5, 0.8, 1.1, 8.0, 3.0,
    ["I gazed \u2014 and gazed \u2014 but little thought",
     "What wealth the show to me had brought"],
    sz=32, color=RGBColor(0xFF,0xFF,0xFF), font='Arial')

accent(s5, 2.8, 4.3, 3.5, 0xFF, 0xD7, 0x00)

multi_tb(s5, 0.8, 4.6, 8.0, 1.2,
    ["\u51DD\u89C6  \u00B7  \u65E0\u5F62\u8D22\u5BCC"],
    sz=34, color=RGBColor(0xFF,0xE0,0x88), font='Microsoft YaHei')

multi_tb(s5, 1.5, 5.3, 7.5, 0.8,
    ["\u8BD7\u4EBA\u4E45\u4E45\u51DD\u89C6\uFF0C\u5F53\u65F6\u672A\u5BDF\u89C9\u8FD9\u662F\u600E\u6837\u7684\u8D22\u5BCC"],
    sz=16, color=RGBColor(0xDD,0xCC,0xBB), font='Microsoft YaHei')

page_no(s5, 5)

# ================================================================
#  SLIDE 6 - Memory: Inward Eye
# ================================================================
s6 = blank_slide()
bg6 = (0x2C, 0x2C, 0x54)
set_bg(s6, *bg6)

for _ in range(25):
    fx = random.uniform(3.0, 10.3)
    fy = random.uniform(2.0, 5.5)
    fr = random.uniform(0.03, 0.15)
    fb = random.randint(200, 255)
    circle(s6, fx, fy, fr, 0xFF, fb // 2, 0x00)

multi_tb(s6, 1.0, 1.2, 11.3, 2.8,
    ["They flash upon that inward eye",
     "the bliss of solitude"],
    sz=38, color=RGBColor(0xFF,0xFF,0xFF), font='Arial')

accent(s6, 5.0, 4.2, 3.3, 0xFF, 0xD7, 0x00)

multi_tb(s6, 1.0, 4.5, 11.3, 1.3,
    ["\u5FC3\u7075\u4E4B\u773C  \u00B7  \u5B64\u5BC2\u4E4B\u798F"],
    sz=34, color=RGBColor(0xFF,0xE0,0xAA), font='Microsoft YaHei')

multi_tb(s6, 1.5, 5.3, 10.3, 0.8,
    ["\u8EBA\u5728\u6C99\u53D1\u4E0A\uFF0C\u7A7A\u865A\u6216\u6C89\u601D\u65F6\uFF0C\u6C34\u4ED9\u5728\u5185\u5FC3\u95EA\u73B0"],
    sz=16, color=RGBColor(0xBB,0xBB,0xEE), font='Microsoft YaHei')

page_no(s6, 6)

# ================================================================
#  SLIDE 7 - Ending: Heart Dances with Flowers
# ================================================================
s7 = blank_slide()
bg7 = (0x1A, 0x1A, 0x2E)
set_bg(s7, *bg7)

circle(s7, 6.66, 3.2, 2.8, *blend(bg7, (0xFF,0xD7,0x00), 0.15))
circle(s7, 6.66, 3.2, 1.8, *blend(bg7, (0xFF,0xC1,0x07), 0.22))
circle(s7, 6.66, 3.2, 0.9, *blend(bg7, (0xFF,0xEE,0x88), 0.35))

multi_tb(s7, 1.0, 1.3, 11.3, 2.8,
    ["And then my heart with pleasure fills,",
     "And dances with the daffodils."],
    sz=36, color=RGBColor(0xFF,0xFF,0xFF), font='Arial')

accent(s7, 5.0, 4.3, 3.3, 0xFF, 0xD7, 0x00)

multi_tb(s7, 1.0, 4.6, 11.3, 1.3,
    ["\u5FC3  \u00B7  \u559C\u60A6  \u00B7  \u540C\u821E"],
    sz=36, color=RGBColor(0xFF,0xE0,0x88), font='Microsoft YaHei')

multi_tb(s7, 2.0, 5.5, 9.3, 0.8,
    ["\u5FC3\u4E2D\u5145\u6EE1\u559C\u60A6\uFF0C\u4E5F\u968F\u6C34\u4ED9\u4E00\u8D77\u8DF3\u821E"],
    sz=16, color=RGBColor(0xCC,0xCC,0xDD), font='Microsoft YaHei')

page_no(s7, 7)

# ================================================================
#  SAVE
# ================================================================
output = r"f:\programming-experiment\I_Wandered_Lonely_as_a_Cloud.pptx"
prs.save(output)
print(f"Done! Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
