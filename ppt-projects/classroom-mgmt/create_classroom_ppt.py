"""Generate Classroom Management System Report PPT based on YIJIAN feedback."""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
TOTAL = 13

# === COLORS ===
DARK_BG = (0x1A, 0x1A, 0x2E)
SECTION_BG = (0x2C, 0x3E, 0x50)
ACCENT_RED = (0xE7, 0x4C, 0x3C)
ACCENT_BLUE = (0x34, 0x98, 0xDB)
ACCENT_GREEN = (0x27, 0xAE, 0x60)
ACCENT_PURPLE = (0x8E, 0x44, 0xAD)
ACCENT_ORANGE = (0xE6, 0x7E, 0x22)
ACCENT_TEAL = (0x1A, 0xBC, 0x9C)
LIGHT_GRAY = (0xEC, 0xF0, 0xF1)
WHITE = (0xFF, 0xFF, 0xFF)
BLACK = (0x00, 0x00, 0x00)
DARK_TEXT = (0x2C, 0x3E, 0x50)
MID_GRAY = (0x99, 0x99, 0x99)
CODE_BG = (0xF5, 0xF5, 0xF5)
CARD_BG = (0xF0, 0xF4, 0xF8)
VERY_LIGHT_BLUE = (0xEB, 0xF5, 0xFB)
VERY_LIGHT_GREEN = (0xE8, 0xF8, 0xF5)
VERY_LIGHT_RED = (0xFD, 0xED, 0xEC)
LIGHT_BORDER = (0xDD, 0xDD, 0xDD)
HEADER_BG = (0xF8, 0xF9, 0xFA)

# === HELPER FUNCTIONS ===

def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, r, g, b):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

def textbox(slide, l, t, w, h, txt, sz=18, bold=False, color=None, align=PP_ALIGN.LEFT, font='Microsoft YaHei', italic=False):
    if color is None:
        color = RGBColor(*BLACK)
    elif isinstance(color, tuple):
        color = RGBColor(*color)
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.font.name = font
    p.font.italic = italic
    return tf

def add_para(tf, text, sz=18, bold=False, level=0, color=None, font='Microsoft YaHei', italic=False, space_after=6):
    if isinstance(color, tuple):
        color = RGBColor(*color)
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.level = level
    p.font.name = font
    p.font.italic = italic
    p.space_after = Pt(space_after)
    return p

def accent_line(slide, l, t, w, r, g, b):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(r, g, b)
    line.line.fill.background()

def dashed_line(slide, l, t, w, r, g, b):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(r, g, b)
    line.line.fill.background()
    line.fill.fore_color.rgb = RGBColor(r, g, b)

def rounded_rect(slide, l, t, w, h, fill_color, border_color=None, corner_radius=0.15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    if border_color:
        shape.line.color.rgb = RGBColor(*border_color)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def rect_shape(slide, l, t, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    if border_color:
        shape.line.color.rgb = RGBColor(*border_color)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_arrow(slide, l, t, w, h, color):
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(*color)
    arrow.line.fill.background()
    return arrow

def oval_shape(slide, l, t, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    shape.line.fill.background()
    return shape

def add_text_to_shape(shape, text, sz=14, bold=False, color=None, align=PP_ALIGN.CENTER, font='Microsoft YaHei'):
    if color is None:
        color = RGBColor(*BLACK)
    elif isinstance(color, tuple):
        color = RGBColor(*color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.font.name = font
    return tf

def set_shape_text_middle(shape):
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    txBody = shape._element.txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')

def page_header(slide, page_num):
    textbox(slide, 0.4, 0.15, 8, 0.4, '教室登记管理系统 - 项目报告',
            sz=9, color=MID_GRAY, font='Microsoft YaHei')
    textbox(slide, 11.5, 7.05, 1.5, 0.35, f'{page_num}/{TOTAL}',
            sz=10, color=MID_GRAY, align=PP_ALIGN.RIGHT, font='Microsoft YaHei')

def section_header(slide, title):
    textbox(slide, 0.6, 0.2, 11.5, 0.8, title, sz=30, bold=True, color=DARK_TEXT)
    accent_line(slide, 0.6, 1.0, 11.5, *ACCENT_RED)

def code_block(slide, l, t, w, h, lines, line_height=0.28, sz=11):
    bg = rounded_rect(slide, l, t, w, h, CODE_BG, LIGHT_BORDER, 0.08)
    tf = textbox(slide, l + 0.15, t + 0.1, w - 0.3, h - 0.2, '', sz=sz, font='Consolas', color=BLACK)
    tf.paragraphs[0].text = lines[0] if lines else ''
    tf.paragraphs[0].font.name = 'Consolas'
    tf.paragraphs[0].font.size = Pt(sz)
    for line in lines[1:]:
        p = add_para(tf, line, sz=sz, font='Consolas', color=BLACK, space_after=2)
    return tf

# ================================================================
#  SLIDE 1 - COVER
# ================================================================
s1 = blank_slide()
set_bg(s1, *DARK_BG)

textbox(s1, 1.0, 1.8, 11.3, 1.2, 'Classroom Management System',
        sz=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Arial')
textbox(s1, 1.0, 2.9, 11.3, 0.9, '教室登记管理系统',
        sz=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Microsoft YaHei')

accent_line(s1, 4.5, 3.8, 4.3, *ACCENT_RED)

textbox(s1, 1.0, 4.1, 11.3, 0.6, '—— 项目报告',
        sz=18, color=(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER, font='Microsoft YaHei', italic=True)

info_tf = textbox(s1, 4.5, 5.2, 5.0, 1.2, '', sz=14, color=(0xBB, 0xBB, 0xBB), align=PP_ALIGN.LEFT, font='Microsoft YaHei')
info_tf.paragraphs[0].text = '姓  名：__________________'
add_para(info_tf, '学  号：__________________', sz=14, color=(0xBB, 0xBB, 0xBB), font='Microsoft YaHei', space_after=4)
add_para(info_tf, '日  期：__________________', sz=14, color=(0xBB, 0xBB, 0xBB), font='Microsoft YaHei', space_after=4)

textbox(s1, 11.5, 7.05, 1.5, 0.35, f'1/{TOTAL}',
        sz=10, color=MID_GRAY, align=PP_ALIGN.RIGHT, font='Microsoft YaHei')

print('Slide 1 (Cover) done.')

# ================================================================
#  SLIDE 2 - PROJECT INTRODUCTION
# ================================================================
s2 = blank_slide()
section_header(s2, '项目简介')
page_header(s2, 2)

sections = [
    ('📋  项目背景', VERY_LIGHT_BLUE, ACCENT_BLUE, [
        '传统教室管理依赖人工登记，效率低下且容易出错',
        '借还教室流程不规范，缺乏统一管理系统',
        '教室资源利用率低，存在空闲浪费现象',
    ]),
    ('🎯  项目目标', VERY_LIGHT_GREEN, ACCENT_GREEN, [
        '开发一套命令行教室登记管理系统',
        '实现教室的借出、归还、预约等核心功能',
        '提供查询统计功能，提升教室管理效率',
    ]),
    ('👥  目标用户', VERY_LIGHT_RED, ACCENT_RED, [
        '学校教务管理人员——日常教室调度与管理',
        '教师与学生——快速查询可用教室并预约',
        '系统管理员——维护系统数据与权限',
    ]),
]

y_start = 1.4
for idx, (title, bg_color, accent_color, bullets) in enumerate(sections):
    y = y_start + idx * 1.85
    rounded_rect(s2, 0.6, y, 12.1, 1.65, bg_color, accent_color, 0.12)
    ov = oval_shape(s2, 0.8, y + 0.15, 0.4, 0.4, accent_color)
    tf = textbox(s2, 1.4, y + 0.05, 11, 0.4, title, sz=19, bold=True, color=accent_color)
    bt = textbox(s2, 1.4, y + 0.55, 11, 1.0, '', sz=14, color=DARK_TEXT)
    for bi, b in enumerate(bullets):
        if bi == 0:
            bt.paragraphs[0].text = '• ' + b
            bt.paragraphs[0].font.size = Pt(14)
            bt.paragraphs[0].font.name = 'Microsoft YaHei'
            bt.paragraphs[0].font.color.rgb = RGBColor(*DARK_TEXT)
        else:
            add_para(bt, '• ' + b, sz=14, color=DARK_TEXT)

print('Slide 2 (Project Intro) done.')

# ================================================================
#  SLIDE 3 - FUNCTION MODULE DIAGRAM
# ================================================================
s3 = blank_slide()
section_header(s3, '功能模块图')
page_header(s3, 3)

root = rounded_rect(s3, 4.8, 1.4, 3.5, 0.8, ACCENT_RED, WHITE, 0.1)
add_text_to_shape(root, '主菜单\nClassroom Management', sz=14, bold=True, color=WHITE)

branches = [
    ('Borrow\n借出教室', (0.6, 2.8)),
    ('Return\n归还教室', (3.3, 2.8)),
    ('View\n查看记录', (6.0, 2.8)),
    ('Query\n查询功能', (8.7, 2.8)),
    ('Reserve\n预约教室', (11.4, 2.8)),
]

for label, (bx, by) in branches:
    box = rounded_rect(s3, bx, by, 2.4, 0.65, ACCENT_BLUE, WHITE, 0.1)
    add_text_to_shape(box, label, sz=12, bold=True, color=WHITE)
    add_arrow(s3, bx + 0.95, by - 0.35, 0.3, 0.35, ACCENT_RED)

sub_items = [
    ('输入教室号/姓名/学号\n自动记录时间', 0.3, 3.85),
    ('输入教室号\n更新状态为空闲', 3.0, 3.85),
    ('显示所有借还记录\n含状态标签', 5.7, 3.85),
    ('按教室/用户/学号\n多维度查询', 8.4, 3.85),
    ('输入时间段\n自动检测冲突', 11.1, 3.85),
]
for txt, sx, sy in sub_items:
    sub = rounded_rect(s3, sx, sy, 3.0, 0.75, LIGHT_GRAY, LIGHT_BORDER, 0.08)
    add_text_to_shape(sub, txt, sz=10, color=DARK_TEXT)

# legend
textbox(s3, 0.6, 5.0, 11.5, 0.4, '说明：主菜单共8个选项，涵盖借出、归还、查看、查询（教室/用户/学号）、预约、退出六大模块',
        sz=10, color=MID_GRAY, align=PP_ALIGN.CENTER)
textbox(s3, 0.6, 5.4, 11.5, 0.4, '蓝色框 = 一级功能模块  |  灰色框 = 功能说明  |  箭头 = 菜单层级关系',
        sz=9, color=MID_GRAY, align=PP_ALIGN.CENTER)

# data flow summary
textbox(s3, 0.6, 5.9, 11.5, 0.8, '数据流向：用户输入 → 菜单路由 → 业务逻辑(registry.c) → 数据操作(classroom.c) → 控制台输出',
        sz=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)

print('Slide 3 (Function Module) done.')

# ================================================================
#  SLIDE 4 - USE CASE & DATA FLOW DIAGRAM
# ================================================================
s4 = blank_slide()
section_header(s4, '用例图 & 数据流图')
page_header(s4, 4)

# LEFT: Use Case Diagram
textbox(s4, 0.6, 1.3, 5.5, 0.5, '用例图 (Use Case Diagram)', sz=18, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

actor = rounded_rect(s4, 0.6, 2.0, 1.8, 0.6, ACCENT_BLUE, WHITE, 0.1)
add_text_to_shape(actor, '👤 用户', sz=14, bold=True, color=WHITE)

use_cases = [
    '借出教室', '归还教室', '查看记录',
    '按教室查询', '按用户查询', '按学号查询', '预约教室',
]
for i, uc in enumerate(use_cases):
    col = i % 3
    row = i // 3
    ux = 3.3 + col * 1.8
    uy = 1.9 + row * 1.0
    if i < 6:
        box = rounded_rect(s4, ux, uy, 1.6, 0.5, VERY_LIGHT_BLUE, ACCENT_BLUE, 0.08)
        add_text_to_shape(box, uc, sz=10, color=DARK_TEXT)
    else:
        box = rounded_rect(s4, 4.2, 3.9, 1.6, 0.5, VERY_LIGHT_GREEN, ACCENT_GREEN, 0.08)
        add_text_to_shape(box, uc, sz=10, color=DARK_TEXT)

# RIGHT: Data Flow Diagram
textbox(s4, 7.0, 1.3, 5.5, 0.5, '顶层数据流图 (DFD Level-0)', sz=18, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

dfd_boxes = [
    ('用户输入', 7.2, 2.0, ACCENT_BLUE),
    ('教室管理系统', 9.0, 2.0, ACCENT_RED),
    ('数据存储\n(Classroom[]/Record[])', 10.8, 2.0, ACCENT_GREEN),
]
for label, dx, dy, col in dfd_boxes:
    box = rounded_rect(s4, dx, dy, 1.6, 1.0, col, WHITE, 0.1)
    add_text_to_shape(box, label, sz=11, bold=True, color=WHITE)

# arrows in DFD
for ax in [8.8, 10.5]:
    arr = add_arrow(s4, ax, 2.7, 0.2, 0.2, MID_GRAY)

textbox(s4, 8.0, 3.3, 4.0, 0.3, '输入/输出 →', sz=9, color=MID_GRAY, align=PP_ALIGN.CENTER)
textbox(s4, 9.7, 3.3, 4.0, 0.3, '读写 →', sz=9, color=MID_GRAY, align=PP_ALIGN.CENTER)

# BOTTOM: Core Data Structures
dt_y = 4.4
textbox(s4, 0.6, dt_y, 11.5, 0.4, '核心数据结构', sz=16, bold=True, color=DARK_TEXT)

# Table
tbl_shape = s4.shapes.add_table(3, 4, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.8))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(1.8)
tbl.columns[1].width = Inches(3.0)
tbl.columns[2].width = Inches(3.0)
tbl.columns[3].width = Inches(4.3)

headers = ['结构体', '类型', '核心字段', '用途']
for ci, hdr in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.text = hdr
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*WHITE)
        p.font.name = 'Microsoft YaHei'
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(*DARK_BG)

data_rows = [
    ['Classroom', 'struct', 'id, name[50], isOccupied, borrower[50], student_id[20]', '存储10间教室的实时状态'],
    ['Record', 'struct', 'classroom[50], borrower[50], student_id[20], isReturned, isReservation, borrowTime, returnTime, reserveStartTime, reserveEndTime', '存储所有借还与预约记录（最多100条）'],
]
for ri, row in enumerate(data_rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = 'Microsoft YaHei'
            if ci == 2:
                p.font.name = 'Consolas'
                p.font.size = Pt(10)
        if ri % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*LIGHT_GRAY)

print('Slide 4 (Use Case & Data Flow) done.')

# ================================================================
#  SLIDE 5 - DATA STRUCTURES
# ================================================================
s5 = blank_slide()
section_header(s5, '数据结构详解')
page_header(s5, 5)

# Classroom struct
textbox(s5, 0.6, 1.3, 5.5, 0.4, '① Classroom 结构体 —— 教室状态', sz=17, bold=True, color=ACCENT_BLUE)
code_block(s5, 0.6, 1.8, 5.8, 2.2, [
    'typedef struct {',
    '    int   id;              // 教室编号 (1-10)',
    '    char  name[50];        // 教室名称',
    '    int   isOccupied;      // 是否被占用 (0/1)',
    '    char  borrower[50];    // 借用人姓名',
    '    char  student_id[20];  // 学号',
    '} Classroom;',
], sz=11)

# Record struct
textbox(s5, 7.0, 1.3, 5.5, 0.4, '② Record 结构体 —— 借还/预约记录', sz=17, bold=True, color=ACCENT_GREEN)
code_block(s5, 7.0, 1.8, 5.8, 2.2, [
    'typedef struct {',
    '    char    classroom[50];      // 教室名',
    '    char    borrower[50];       // 借用人',
    '    char    student_id[20];     // 学号',
    '    int     isReturned;         // 是否已还',
    '    int     isReservation;      // 是否预约',
    '    time_t  borrowTime;         // 借出时间',
    '    time_t  returnTime;         // 归还时间',
    '    time_t  reserveStartTime;   // 预约开始',
    '    time_t  reserveEndTime;     // 预约结束',
    '} Record;',
], sz=11)

# Design points
textbox(s5, 0.6, 4.3, 11.5, 0.4, '设计要点', sz=17, bold=True, color=DARK_TEXT)
design_points = [
    'isReservation 标志位区分"借还记录"与"预约记录"，两者共用同一 Record 结构体',
    'time_t 类型统一存储时间戳，支持跨天、跨月的时间比较与格式化输出',
    '全局数组 classrooms[10] + records[100] 实现内存级数据管理，无需外部数据库',
    '所有字符串字段使用固定长度 char[] 而非指针，避免动态内存管理的复杂性',
]
dt = textbox(s5, 0.8, 4.8, 11.5, 1.8, '', sz=14, color=DARK_TEXT)
for di, dp in enumerate(design_points):
    if di == 0:
        dt.paragraphs[0].text = f'• {dp}'
        dt.paragraphs[0].font.size = Pt(14)
        dt.paragraphs[0].font.name = 'Microsoft YaHei'
        dt.paragraphs[0].font.color.rgb = RGBColor(*DARK_TEXT)
    else:
        add_para(dt, f'• {dp}', sz=14, color=DARK_TEXT)

print('Slide 5 (Data Structures) done.')

# ================================================================
#  SLIDE 6 - CORE FUNCTIONS
# ================================================================
s6 = blank_slide()
section_header(s6, '核心函数')
page_header(s6, 6)

functions = [
    ('borrowClassroom()', '借出教室：选择教室号 → 输入姓名学号 → 自动记录借出时间', ACCENT_BLUE),
    ('returnClassroom()', '归还教室：选择教室号 → 自动标记已归还 → 记录归还时间', ACCENT_GREEN),
    ('reserveClassroom()', '预约教室：输入时间段 → 校验合法性 → 检测冲突 → 创建预约记录', ACCENT_RED),
    ('addRecord()', '添加记录：将借出信息写入 records[] 数组，更新教室状态为已占用', ACCENT_PURPLE),
    ('displayAllRecords()', '查看记录：遍历所有记录，区分借还/预约类型，显示状态标签', ACCENT_ORANGE),
    ('queryByStudentID()', '学号查询：按学号检索所有关联记录，支持模糊匹配', ACCENT_TEAL),
]

for i, (name, desc, col) in enumerate(functions):
    row = i // 2
    ccol = i % 2
    fx = 0.6 + ccol * 6.2
    fy = 1.4 + row * 1.75

    card = rounded_rect(s6, fx, fy, 5.8, 1.45, VERY_LIGHT_BLUE if ccol == 0 else VERY_LIGHT_GREEN, col, 0.12)

    # accent bar on left
    bar = rect_shape(s6, fx, fy, 0.08, 1.45, col)

    name_tf = textbox(s6, fx + 0.25, fy + 0.1, 5.3, 0.5, name, sz=16, bold=True, color=col, font='Consolas')
    desc_tf = textbox(s6, fx + 0.25, fy + 0.65, 5.3, 0.65, desc, sz=12, color=DARK_TEXT)

print('Slide 6 (Core Functions) done.')

# ================================================================
#  SLIDE 7 - KEY DIFFICULTIES (1/2)
# ================================================================
s7 = blank_slide()
section_header(s7, '关键难点 (1/2)')
page_header(s7, 7)

difficulties_left = [
    {
        'title': '冲突检测',
        'problem': '预约教室时需确保时间段不与已有预约重叠',
        'solution': '使用区间重叠算法：检查新区间 [S,E] 是否与已有区间 [s,e] 相交',
        'insight': '核心逻辑：!(E < s || S > e) 判断两区间是否有交集',
        'code': ['int hasReservationConflict(', '    const char* classroom,', '    time_t startTime,', '    time_t endTime)', '{', '  for (i=0; i<recordCount; i++) {', '    if (!(endTime < records[i].', '        reserveStartTime ||', '        startTime > records[i].', '        reserveEndTime))', '      return 1; // 冲突', '  }', '  return 0; // 无冲突', '}'],
    },
    {
        'title': '过去时间限制',
        'problem': '用户可能输入过去的日期时间，导致无效预约',
        'solution': '在解析时间后与当前时间 time(NULL) 比较，拒绝过去的时间',
        'insight': '简单但关键：`if (startTime <= now) return error;` 防止历史预约',
        'code': ['now = time(NULL);', 'if (startTime <= now) {', '  printMessageBox(', '    "Start time must be', '     in the future.");', '  return;', '}'],
    },
]

for idx, diff in enumerate(difficulties_left):
    dx = 0.4 + idx * 6.3
    dy = 1.4

    # Title
    textbox(s7, dx, dy, 5.8, 0.4, f'难点 {idx+1}：{diff["title"]}', sz=18, bold=True, color=ACCENT_RED)

    # Problem block
    pb = rounded_rect(s7, dx, dy + 0.45, 5.8, 0.55, VERY_LIGHT_RED, ACCENT_RED, 0.08)
    add_text_to_shape(pb, '🔴 Problem：' + diff['problem'], sz=10, color=DARK_TEXT, align=PP_ALIGN.LEFT)
    pb.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'

    # Solution block
    sb = rounded_rect(s7, dx, dy + 1.05, 5.8, 0.55, VERY_LIGHT_GREEN, ACCENT_GREEN, 0.08)
    add_text_to_shape(sb, '🟢 Solution：' + diff['solution'], sz=10, color=DARK_TEXT, align=PP_ALIGN.LEFT)
    sb.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'

    # Key Insight block
    ib = rounded_rect(s7, dx, dy + 1.65, 5.8, 0.55, VERY_LIGHT_BLUE, ACCENT_BLUE, 0.08)
    add_text_to_shape(ib, '💡 Key Insight：' + diff['insight'], sz=10, color=DARK_TEXT, align=PP_ALIGN.LEFT)
    ib.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'

    # Code snippet
    code_block(s7, dx, dy + 2.35, 5.8, 2.8, diff['code'], sz=9)

print('Slide 7 (Key Difficulties 1/2) done.')

# ================================================================
#  SLIDE 8 - KEY DIFFICULTIES (2/2)
# ================================================================
s8 = blank_slide()
section_header(s8, '关键难点 (2/2)')
page_header(s8, 8)

# Three smaller difficulty boxes
diffs_right = [
    {
        'num': '3',
        'title': '时间格式解析',
        'desc': '用户输入 "MM-DD-HH" 格式，需手动解析为 time_t。使用 sscanf 拆分月/日/时，结合当前年份构造 tm 结构体，调用 mktime() 转换。注意处理闰年、跨月等边界情况。',
        'color': ACCENT_PURPLE,
    },
    {
        'num': '4',
        'title': '预约与借还的互斥',
        'desc': '当前时间处于预约时段内的教室不可被借出。borrowClassroom() 中先调用 isRoomReservedAtTime() 检查，若教室已被预约则拒绝借出操作，确保预约优先级。',
        'color': ACCENT_ORANGE,
    },
    {
        'num': '5',
        'title': '归还记录的匹配',
        'desc': '归还时需找到最近一条"未归还的借出记录"来更新状态。returnClassroomRecord() 从 recordCount-1 逆序遍历，匹配教室名 + 非预约 + 未归还三个条件，确保更新正确的记录。',
        'color': ACCENT_TEAL,
    },
]

for i, d in enumerate(diffs_right):
    dx = 0.4 + i * 4.2
    dy = 1.3

    card = rounded_rect(s8, dx, dy, 4.0, 2.5, VERY_LIGHT_BLUE if i % 2 == 0 else VERY_LIGHT_GREEN, d['color'], 0.12)
    ov = oval_shape(s8, dx + 1.5, dy + 0.1, 0.5, 0.5, d['color'])
    add_text_to_shape(ov, d['num'], sz=16, bold=True, color=WHITE)
    textbox(s8, dx + 0.2, dy + 0.75, 3.6, 0.4, d['title'], sz=15, bold=True, color=d['color'])
    textbox(s8, dx + 0.2, dy + 1.2, 3.6, 1.2, d['desc'], sz=10, color=DARK_TEXT)

# State machine table
textbox(s8, 0.6, 4.1, 11.5, 0.4, '预约状态机（now 与 reserveStart/End 比较）', sz=15, bold=True, color=DARK_TEXT)

tbl_shape2 = s8.shapes.add_table(3, 4, Inches(2.5), Inches(4.55), Inches(8.3), Inches(1.2))
tbl2 = tbl_shape2.table
tbl2.columns[0].width = Inches(2.0)
tbl2.columns[1].width = Inches(2.1)
tbl2.columns[2].width = Inches(2.1)
tbl2.columns[3].width = Inches(2.1)

h2 = ['条件', 'now < startTime', 'startTime ≤ now ≤ endTime', 'now > endTime']
for ci, h in enumerate(h2):
    cell = tbl2.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*WHITE)
        p.font.name = 'Microsoft YaHei'
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(*DARK_BG)

r2 = [['状态', 'Upcoming (即将到来)', 'In Progress (进行中)', 'Expired (已过期)'],
      ['含义', '预约尚未开始', '当前处于预约时段内', '预约已结束']]
for ri, row in enumerate(r2):
    for ci, val in enumerate(row):
        cell = tbl2.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = 'Microsoft YaHei'
            if ri == 0:
                p.font.bold = True

# Summary row
textbox(s8, 0.6, 6.1, 11.5, 0.6, '🔑 难点小结：五大难点涵盖 时间处理、冲突检测、状态管理、输入解析、数据匹配 五个核心领域，'
        '体现了命令行系统设计的典型挑战与工程实践。',
        sz=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)

print('Slide 8 (Key Difficulties 2/2) done.')

# ================================================================
#  SLIDE 9 - SCREENSHOTS
# ================================================================
s9 = blank_slide()
section_header(s9, '运行截图')
page_header(s9, 9)

screenshots = [
    '主菜单界面', '借出教室操作', '归还教室操作',
    '查看所有记录', '按教室查询', '预约教室流程',
]

for i, label in enumerate(screenshots):
    col = i % 3
    row = i // 3
    sx = 0.6 + col * 4.2
    sy = 1.5 + row * 2.7

    placeholder = rounded_rect(s9, sx, sy, 3.8, 2.3, CODE_BG, LIGHT_BORDER, 0.15)
    add_text_to_shape(placeholder, f'📷 {label}\n\n（待补充实际截图）', sz=13, color=MID_GRAY)
    placeholder.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'
    set_shape_text_middle(placeholder)

print('Slide 9 (Screenshots) done.')

# ================================================================
#  SLIDE 10 - DEMO FLOW
# ================================================================
s10 = blank_slide()
section_header(s10, '演示流程')
page_header(s10, 10)

steps = [
    ('①', '启动程序', '运行程序，显示主菜单界面，8个功能选项清晰列出'),
    ('②', '借出教室', '选择 Borrow(1) → 输入教室号 → 输入姓名学号 → 系统记录时间'),
    ('③', '查看记录', '选择 View Records(3) → 显示所有教室状态与借还历史'),
    ('④', '预约教室', '选择 Reserve(7) → 输入教室号 → 输入起止时间 → 冲突检测'),
    ('⑤', '查询操作', '选择 Query(4/5/6) → 按教室/用户/学号查询 → 显示结果'),
    ('⑥', '归还教室', '选择 Return(2) → 输入教室号 → 自动更新状态与归还时间'),
]

for i, (num, title, desc) in enumerate(steps):
    sy = 1.5 + i * 0.95

    ov = oval_shape(s10, 0.8, sy + 0.05, 0.55, 0.55, ACCENT_BLUE)
    add_text_to_shape(ov, num, sz=16, bold=True, color=WHITE)

    textbox(s10, 1.6, sy + 0.05, 3.0, 0.5, title, sz=15, bold=True, color=DARK_TEXT)
    textbox(s10, 4.6, sy + 0.05, 8.0, 0.5, desc, sz=12, color=DARK_TEXT)

    if i < len(steps) - 1:
        add_arrow(s10, 1.0, sy + 0.7, 0.15, 0.25, ACCENT_BLUE)

print('Slide 10 (Demo Flow) done.')

# ================================================================
#  SLIDE 11 - FEATURES & ADVANTAGES
# ================================================================
s11 = blank_slide()
section_header(s11, '特色与优势')
page_header(s11, 11)

textbox(s11, 0.6, 1.2, 11.5, 0.5, '本系统具备以下 10 项核心优势', sz=16, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

features = [
    ('⏱️  实时时间记录', '自动调用 time() 记录操作时间'),
    ('🔍  多维查询', '按教室/用户/学号精准检索'),
    ('📅  预约管理', '支持预约时段，自动冲突检测'),
    ('🛡️  输入校验', '多层输入验证，防止非法操作'),
    ('📊  状态可视化', '清晰展示教室占用/空闲/预约状态'),
    ('🔄  完整生命周期', '借出→归还 全流程闭环管理'),
    ('🧩  模块化设计', 'registry.c + classroom.c 职责分离'),
    ('💾  内存数据管理', '无需外部数据库，轻量高效'),
    ('📝  结构化记录', 'Record 结构体统一管理借还与预约'),
    ('🚦  智能状态机', 'Upcoming/In Progress/Expired 三态'),
]

for i, (name, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    fx = 0.6 + col * 6.2
    fy = 1.85 + row * 0.95

    card = rounded_rect(s11, fx, fy, 5.8, 0.8, VERY_LIGHT_BLUE if col == 0 else VERY_LIGHT_GREEN, ACCENT_BLUE if col == 0 else ACCENT_GREEN, 0.1)
    textbox(s11, fx + 0.2, fy + 0.05, 5.4, 0.35, name, sz=13, bold=True, color=DARK_TEXT)
    textbox(s11, fx + 0.2, fy + 0.4, 5.4, 0.35, desc, sz=10, color=MID_GRAY)

print('Slide 11 (Features & Advantages) done.')

# ================================================================
#  SLIDE 12 - SUMMARY
# ================================================================
s12 = blank_slide()
section_header(s12, '项目总结')
page_header(s12, 12)

columns = [
    {
        'emoji': '✅',
        'title': '已完成',
        'color': ACCENT_GREEN,
        'bg': VERY_LIGHT_GREEN,
        'items': [
            '教室借出与归还功能',
            '多维度查询系统',
            '预约与冲突检测',
            '命令行交互界面',
            '完整记录追踪',
        ],
    },
    {
        'emoji': '📚',
        'title': '经验教训',
        'color': ACCENT_BLUE,
        'bg': VERY_LIGHT_BLUE,
        'items': [
            'C语言结构体设计',
            '时间处理与比较',
            '区间重叠算法',
            '状态机设计模式',
            '输入校验最佳实践',
        ],
    },
    {
        'emoji': '🚀',
        'title': '未来改进',
        'color': ACCENT_RED,
        'bg': VERY_LIGHT_RED,
        'items': [
            'GUI 图形界面',
            '数据库持久化存储',
            '用户认证与权限管理',
            'Web 远程访问',
            '数据统计分析报表',
        ],
    },
]

for ci, col_data in enumerate(columns):
    cx = 0.6 + ci * 4.2
    cy = 1.5

    header_bg = rounded_rect(s12, cx, cy, 3.8, 0.6, col_data['color'], WHITE, 0.08)
    add_text_to_shape(header_bg, f'{col_data["emoji"]}  {col_data["title"]}', sz=18, bold=True, color=WHITE)

    card = rounded_rect(s12, cx, cy + 0.6, 3.8, 3.0, col_data['bg'], col_data['color'], 0.1)

    item_tf = textbox(s12, cx + 0.2, cy + 0.75, 3.4, 2.5, '', sz=12, color=DARK_TEXT)
    for ii, item in enumerate(col_data['items']):
        if ii == 0:
            item_tf.paragraphs[0].text = f'• {item}'
            item_tf.paragraphs[0].font.size = Pt(13)
            item_tf.paragraphs[0].font.name = 'Microsoft YaHei'
            item_tf.paragraphs[0].font.color.rgb = RGBColor(*DARK_TEXT)
        else:
            add_para(item_tf, f'• {item}', sz=13, color=DARK_TEXT)

    if ci == 2:
        dashed_line(s12, cx + 0.3, cy + 3.7, 3.2, *ACCENT_RED)

textbox(s12, 0.6, 5.5, 11.5, 0.5, '--- 虚线表示未来改进为开放项，欢迎讨论 ---',
        sz=10, color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)

print('Slide 12 (Summary) done.')

# ================================================================
#  SLIDE 13 - THANKS
# ================================================================
s13 = blank_slide()
set_bg(s13, *DARK_BG)

textbox(s13, 1.0, 2.0, 11.3, 1.5, 'Thank You!',
        sz=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Arial')
textbox(s13, 1.0, 3.3, 11.3, 1.0, '感谢聆听！',
        sz=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Microsoft YaHei')

accent_line(s13, 4.5, 4.5, 4.3, *ACCENT_RED)

textbox(s13, 1.0, 4.9, 11.3, 0.6, 'Q & A',
        sz=22, color=(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER, font='Arial')

textbox(s13, 1.0, 5.8, 11.3, 0.4, '教室登记管理系统 - 项目报告',
        sz=12, color=MID_GRAY, align=PP_ALIGN.CENTER, font='Microsoft YaHei')

textbox(s13, 11.5, 7.05, 1.5, 0.35, f'13/{TOTAL}',
        sz=10, color=MID_GRAY, align=PP_ALIGN.RIGHT, font='Microsoft YaHei')

print('Slide 13 (Thanks) done.')

# ================================================================
#  SAVE
# ================================================================
output_path = r'f:\programming-experiment\Classroom_Management_Report.pptx'
prs.save(output_path)
print(f'\nDone! Saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
